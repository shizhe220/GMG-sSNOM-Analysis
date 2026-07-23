"""Build the Hankel-fit + FFT overview deck for graphene_4x1_manual lp2 (interior),
wn=860-1000cm-1: one slide per wn, title on the left, the real-space fit comparison
plot top-right, the FFT (spatial data + amp/phase/complex) 4-panel figure below
spanning most of the width.

860-890cm-1 (single-wave regime): fit plot = compare_cavity_models's Hankel + 1/sqrt(x)
"sorted by AIC" comparison ({wn}cm-1_singlewave.png).
900-1000cm-1 (two-wave regime): fit plot = single-wave-vs-two-wave Hankel comparison
({wn}cm-1_hankel2wave.png), same convention as lp1's own deck
(slides/GMG_two_wave_hankel_900to1000.pptx).

Images extracted from fitting_pipeline_graphene_4x1_lp2.ipynb's already-executed cell
outputs by scripts/extract_hankel_fft_slidefigs_lp2.py (see
figures/graphene_4x1_manual/lp2/two_wave_hankel/) -- this script only lays them out, it
doesn't regenerate the fits.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

os.chdir(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

FIG_DIR = 'figures/graphene_4x1_manual/lp2/two_wave_hankel'
OUT_PPTX = 'slides/GMG_hankel_fft_overview_graphene_4x1_lp2.pptx'
FONT = 'Aptos'

NAVY = RGBColor(0x1E, 0x27, 0x61)
MUTED = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

single_wns = [860, 870, 880, 890]
two_wns = [900, 911, 920, 930, 941, 950, 960, 970, 980, 991, 1000]


def img_size_in(path):
    im = Image.open(path)
    return im.size


def no_autofit(tf):
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        el = bodyPr.find(qn(tag))
        if el is not None:
            bodyPr.remove(el)
    bodyPr.append(bodyPr.makeelement(qn('a:noAutofit'), {}))


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             italic=False, font=FONT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ---------------------------------------------------------------- title slide
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.3),
         "Hankel Fit + FFT Overview", 38, WHITE, bold=True, font=FONT)
add_text(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.6),
         "graphene_4x1_manual lp2 (interior), wn = 860-1000 cm⁻¹", 18, RGBColor(0xCA, 0xDC, 0xFC), font=FONT)
add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.1),
         "860-890 cm⁻¹: single-wave Hankel vs. 1/√x real-space comparison. "
         "900-1000 cm⁻¹: single-wave vs. two-wave (edge-launched q + tip-launched 2q) "
         "Hankel fit, both against the FFT's own 2-Lorentzian-peak decomposition.", 13,
         RGBColor(0xCA, 0xDC, 0xFC), italic=True, font=FONT)

# ---------------------------------------------------------------- per-wn slides
margin = Inches(0.35)
title_w = Inches(2.0)


def add_wn_slide(wn, fit_path):
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

    add_text(s, margin, Inches(0.3), title_w, Inches(1.0),
             f"{wn}cm⁻¹", 34, NAVY, bold=True, font=FONT)

    fft_path = f"{FIG_DIR}/{wn}cm-1_fft.png"

    # Fit comparison plot -- top right, height-budgeted (roughly square/portrait, sizing
    # by leftover width would make it far too tall) and centered in the space right of
    # the title column.
    fit_top_y = Inches(0.15)
    fit_budget_h = Inches(3.9)
    pw, ph = img_size_in(fit_path)
    fit_h = fit_budget_h
    fit_w = Emu(int(fit_h * pw / ph))
    avail_left = margin + title_w
    avail_right = SLIDE_W - margin
    fit_x = Emu(int(avail_left + (avail_right - avail_left - fit_w) / 2))
    s.shapes.add_picture(fit_path, fit_x, fit_top_y, width=fit_w, height=fit_h)

    # FFT 4-panel -- below, full content width, height-capped to whatever's left so it
    # never overruns the slide (this figure is wide/short, ~4:1, so width-driven sizing
    # is usually fine, but cap it defensively anyway)
    fft_y = fit_top_y + fit_h + Inches(0.15)
    fft_w = SLIDE_W - 2 * margin
    fw_px, fh_px = img_size_in(fft_path)
    fft_h = Emu(int(fft_w * fh_px / fw_px))
    max_fft_h = SLIDE_H - fft_y - Inches(0.15)
    if fft_h > max_fft_h:
        fft_h = max_fft_h
        fft_w = Emu(int(fft_h * fw_px / fh_px))
    fft_x = Emu(int((SLIDE_W - fft_w) / 2))
    s.shapes.add_picture(fft_path, fft_x, fft_y, width=fft_w, height=fft_h)


for wn in single_wns:
    add_wn_slide(wn, f"{FIG_DIR}/{wn}cm-1_singlewave.png")
for wn in two_wns:
    add_wn_slide(wn, f"{FIG_DIR}/{wn}cm-1_hankel2wave.png")

prs.save(OUT_PPTX)
print(f"Saved {OUT_PPTX}")
