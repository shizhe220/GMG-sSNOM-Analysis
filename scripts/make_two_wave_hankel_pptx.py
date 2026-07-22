"""Build the two-wave Hankel (q + 2q) overview deck for graphene_4x1_manual lp1,
wn=900-1000cm-1: one slide per wn, title on the left, the single-wave-vs-two-wave
Hankel fit comparison plot top-right, the FFT (spatial data + amp/phase/complex)
4-panel figure below spanning most of the width.

Images extracted from fitting_pipeline_graphene_4x1_lp1.ipynb's already-executed
cell outputs (see figures/graphene_4x1_manual/lp1/two_wave_hankel/) -- this script
only lays them out, it doesn't regenerate the fits.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

os.chdir(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

FIG_DIR = 'figures/graphene_4x1_manual/lp1/two_wave_hankel'
OUT_PPTX = 'slides/GMG_two_wave_hankel_900to1000.pptx'
FONT = 'Aptos'

NAVY = RGBColor(0x1E, 0x27, 0x61)
MUTED = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

wns = [900, 911, 920, 930, 941, 950, 960, 970, 980, 991, 1000]


def img_size_in(path):
    im = Image.open(path)
    return im.size


def no_autofit(tf):
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        el = bodyPr.find(qn(tag))
        if el is not None:
            bodyPr.remove(el)
    bodyPr.append(bodyPr.makeelement(qn('a:noAutofit'), {}))


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             italic=False, font=FONT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ---------------------------------------------------------------- title slide
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.3),
         "Two-wave Hankel Fit (q + 2q)", 38, WHITE, bold=True, font=FONT)
add_text(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.6),
         "graphene_4x1_manual lp1 (edge), wn = 900-1000 cm⁻¹", 18, RGBColor(0xCA, 0xDC, 0xFC), font=FONT)
add_text(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.9),
         "Single-wave Hankel vs. two-wave (edge-launched q + tip-launched 2q) fit, "
         "compared against the FFT's own 2-Lorentzian-peak decomposition", 13,
         RGBColor(0xCA, 0xDC, 0xFC), italic=True, font=FONT)

# ---------------------------------------------------------------- per-wn slides
margin = Inches(0.35)
title_w = Inches(2.0)

for wn in wns:
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

    add_text(s, margin, Inches(0.3), title_w, Inches(1.0),
             f"{wn}cm⁻¹", 34, NAVY, bold=True, font=FONT)

    hankel_path = f"{FIG_DIR}/{wn}cm-1_hankel2wave.png"
    fft_path = f"{FIG_DIR}/{wn}cm-1_fft_slidefont.png"

    # Hankel fit comparison -- top right, height-budgeted (this plot is roughly square,
    # sizing it by the leftover width would make it far too tall for the slide) and
    # centered in the space right of the title column.
    hankel_top_y = Inches(0.15)
    hankel_budget_h = Inches(3.9)
    pw, ph = img_size_in(hankel_path)
    hankel_h = hankel_budget_h
    hankel_w = Emu(int(hankel_h * pw / ph))
    avail_left = margin + title_w
    avail_right = SLIDE_W - margin
    hankel_x = Emu(int(avail_left + (avail_right - avail_left - hankel_w) / 2))
    s.shapes.add_picture(hankel_path, hankel_x, hankel_top_y, width=hankel_w, height=hankel_h)

    # FFT 4-panel -- below, full content width, height-capped to whatever's left so it
    # never overruns the slide (this figure is wide/short, ~4:1, so width-driven sizing
    # is usually fine, but cap it defensively anyway)
    fft_y = hankel_top_y + hankel_h + Inches(0.15)
    fft_w = SLIDE_W - 2 * margin
    fw_px, fh_px = img_size_in(fft_path)
    fft_h = Emu(int(fft_w * fh_px / fw_px))
    max_fft_h = SLIDE_H - fft_y - Inches(0.15)
    if fft_h > max_fft_h:
        fft_h = max_fft_h
        fft_w = Emu(int(fft_h * fw_px / fh_px))
    fft_x = Emu(int((SLIDE_W - fft_w) / 2))
    s.shapes.add_picture(fft_path, fft_x, fft_y, width=fft_w, height=fft_h)

prs.save(OUT_PPTX)
print(f"Saved {OUT_PPTX}")
