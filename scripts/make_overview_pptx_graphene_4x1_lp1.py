"""Build the CHT/Hankel/1-sqrtx/FFT overview deck for the graphene_4x1_manual
lp1 (left edge) dataset: one title slide, one slide per wavenumber (CHT figure
on top, real-space + FFT figures below), and a closing summary table slide.

Sibling of scripts/make_overview_pptx_graphene_3x1_manual.py, pointed at
graphene_4x1_manual/lp1."""
import csv
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

os.chdir(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

FIG_DIR = 'figures/graphene_4x1_manual/lp1'
CSV_PATH = 'data/cht_vs_realspace_wavelength_comparison_graphene_4x1_lp1.csv'
OUT_PPTX = 'slides/GMG_CHT_overview_graphene_4x1_lp1.pptx'

NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
TEAL = RGBColor(0x1C, 0x72, 0x93)
MUTED = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_TINT = RGBColor(0xEE, 0xF4, 0xFB)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

wns = [1000, 991, 980, 970, 960, 950, 941, 930, 920, 911, 900, 890, 880, 870, 860]

rows = {}
with open(CSV_PATH) as f:
    for r in csv.DictReader(f):
        rows[int(r['wn_cm1'])] = r


def img_size_in(path):
    im = Image.open(path)
    return im.size


def no_autofit(tf):
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        el = bodyPr.find(qn(tag))
        if el is not None:
            bodyPr.remove(el)
    bodyPr.append(bodyPr.makeelement(qn('a:noAutofit'), {}))


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             italic=False, font='Calibri', anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ---------------------------------------------------------------- title slide
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
bg.shadow.inherit = False

for cx, cy, d, col in [(11.6, -0.9, 4.2, TEAL), (-1.2, 6.6, 3.6, TEAL)]:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = col
    c.fill.transparency = 0
    c.line.fill.background()
    c.shadow.inherit = False
    sp = c.fill.fore_color._xFill
    alpha = sp.makeelement(qn('a:alpha'), {'val': '18000'})
    srgb = sp.find(qn('a:srgbClr'))
    srgb.append(alpha)

add_text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.5),
         "GMG Graphene Plasmon Fitting", 40, WHITE, bold=True, font='Cambria')
add_text(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.7),
         "CHT vs. Real-Space (Hankel / 1/√x) vs. FFT — overview across 15 wavenumbers", 18, ICE)
add_text(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(0.5),
         "860–1000 cm⁻¹  ·  nano-FTIR amplitude line profiles  ·  graphene_4x1_manual dataset "
         "(lp1, left edge — manually-extracted+aligned linecuts from processed_4x1um/860-1000_AVG.npz)",
         13, ICE, italic=True)

# ---------------------------------------------------------------- per-wn slides
margin = Inches(0.5)
content_w = SLIDE_W - 2 * margin

for wn in wns:
    r = rows[wn]
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

    add_text(s, margin, Inches(0.25), content_w, Inches(0.55),
             f"{wn} cm⁻¹", 30, NAVY, bold=True, font='Cambria')

    caption = (f"CHT λₚ = {r['lambda_cht_nm']} nm    "
               f"Hankel λₚ = {r['lambda_hankel_nm']} nm    "
               f"1/√x λₚ = {r['lambda_sqrtx_nm']} nm        "
               f"Δ vs Hankel = {r['cht_vs_hankel_pct']}%    "
               f"Δ vs 1/√x = {r['cht_vs_sqrtx_pct']}%")
    add_text(s, margin, Inches(0.83), content_w, Inches(0.4), caption, 14, MUTED)

    cht_path = f"{FIG_DIR}/cht/{wn}cm-1_cht.png"
    pw, ph = img_size_in(cht_path)
    cht_h = content_w * ph / pw
    cht_y = Inches(1.32)
    s.shapes.add_picture(cht_path, margin, cht_y, width=content_w, height=cht_h)

    rs_path = f"{FIG_DIR}/realspace/{wn}cm-1_realspace.png"
    fft_path = f"{FIG_DIR}/fft/{wn}cm-1_fft.png"
    rs_w_px, rs_h_px = img_size_in(rs_path)
    fft_w_px, fft_h_px = img_size_in(fft_path)

    bottom_h = Inches(1.95)
    rs_w = Emu(int(bottom_h * rs_w_px / rs_h_px))
    fft_w = Emu(int(bottom_h * fft_w_px / fft_h_px))
    gap = Inches(0.2)
    total_w = rs_w + gap + fft_w
    left_x = Emu(int((SLIDE_W - total_w) / 2))
    bottom_y = cht_y + Emu(int(cht_h)) + Inches(0.12)

    s.shapes.add_picture(rs_path, left_x, bottom_y, width=rs_w, height=bottom_h)
    s.shapes.add_picture(fft_path, Emu(int(left_x + rs_w + gap)), bottom_y, width=fft_w, height=bottom_h)

# ---------------------------------------------------------------- summary table slide
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
add_text(s, margin, Inches(0.3), content_w, Inches(0.55), "Summary: λₚ Across All Wavenumbers",
         26, NAVY, bold=True, font='Cambria')

headers = ["wn (cm⁻¹)", "λ CHT (nm)", "λ Hankel (nm)", "λ 1/√x (nm)",
           "k_fit_range", "Δ vs Hankel", "Δ vs 1/√x"]
n_rows = len(wns) + 1
n_cols = len(headers)
table_h = Inches(6.4)
tbl_shape = s.shapes.add_table(n_rows, n_cols, margin, Inches(1.0), content_w, table_h)
tbl = tbl_shape.table
col_w = [Inches(1.5), Inches(1.6), Inches(1.7), Inches(1.6), Inches(2.0), Inches(1.7), Inches(1.7)]
for i, w in enumerate(col_w):
    tbl.columns[i].width = w

for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(13)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

for i, wn in enumerate(wns):
    r = rows[wn]
    vals = [str(wn), r['lambda_cht_nm'], r['lambda_hankel_nm'], r['lambda_sqrtx_nm'],
            r['k_fit_range_1e5cm1'], f"{r['cht_vs_hankel_pct']}%", f"{r['cht_vs_sqrtx_pct']}%"]
    for j, v in enumerate(vals):
        cell = tbl.cell(i + 1, j)
        cell.text = v
        cell.fill.solid()
        cell.fill.fore_color.rgb = ROW_TINT if i % 2 == 0 else WHITE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = MUTED
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

prs.save(OUT_PPTX)
print(f"Saved {OUT_PPTX}")
